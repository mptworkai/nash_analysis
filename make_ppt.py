#!/usr/bin/env python3
"""
Generate a PowerPoint overview of the Nash Analysis tool.

Usage:
  python make_ppt.py                  # writes nash_analysis_overview.pptx
  python make_ppt.py --out FILE.pptx
  python make_ppt.py --help
"""

import sys

if "-h" in sys.argv or "--help" in sys.argv:
    print(__doc__)
    sys.exit(0)

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt
except ImportError:
    print("Missing dependency. Run: pip install python-pptx", file=sys.stderr)
    sys.exit(1)

from pathlib import Path

out_path = Path("nash_analysis_overview.pptx")
args = sys.argv[1:]
i = 0
while i < len(args):
    if args[i] == "--out" and i + 1 < len(args):
        out_path = Path(args[i + 1]); i += 2
    else:
        i += 1

# ── palette ───────────────────────────────────────────────────────────────────
NAVY   = RGBColor(0x1E, 0x3A, 0x5F)
GOLD   = RGBColor(0xC8, 0x9B, 0x2E)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY  = RGBColor(0xF3, 0xF4, 0xF6)
DGRAY  = RGBColor(0x6B, 0x72, 0x80)
GREEN  = RGBColor(0x06, 0x5F, 0x46)
RED    = RGBColor(0x7F, 0x1D, 0x1D)
LTBLUE = RGBColor(0xDB, 0xEA, 0xFE)

W = Inches(13.333)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]  # completely blank


# ── helpers ───────────────────────────────────────────────────────────────────

def add_rect(slide, x, y, w, h, fill=None, line=None):
    shape = slide.shapes.add_shape(1, x, y, w, h)  # MSO_SHAPE_TYPE.RECTANGLE
    shape.line.fill.background() if line is None else None
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line is None:
        shape.line.fill.background()
    return shape


def add_text(slide, text, x, y, w, h, size=18, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb


def bullet_box(slide, items, x, y, w, h, size=16, color=NAVY,
               title=None, title_color=NAVY):
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf  = txb.text_frame
    tf.word_wrap = True
    first = True
    if title:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        run = p.add_run()
        run.text = title
        run.font.size  = Pt(size + 2)
        run.font.bold  = True
        run.font.color.rgb = title_color
    for item in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_before = Pt(4)
        run = p.add_run()
        run.text = item
        run.font.size  = Pt(size)
        run.font.color.rgb = color
    return txb


def navy_header(slide, title, subtitle=None):
    add_rect(slide, 0, 0, W, Inches(1.3), fill=NAVY)
    add_text(slide, title, Inches(0.5), Inches(0.15), Inches(12), Inches(0.75),
             size=32, bold=True, color=WHITE)
    if subtitle:
        add_text(slide, subtitle, Inches(0.5), Inches(0.85), Inches(12), Inches(0.4),
                 size=16, color=GOLD)
    add_rect(slide, 0, Inches(1.3), W, Inches(0.06), fill=GOLD)


def card(slide, x, y, w, h, title, body_lines, title_bg=NAVY, body_bg=LGRAY,
         title_size=16, body_size=14):
    add_rect(slide, x, y, w, Inches(0.45), fill=title_bg)
    add_text(slide, title, x + Inches(0.15), y + Inches(0.05),
             w - Inches(0.3), Inches(0.38),
             size=title_size, bold=True, color=WHITE)
    add_rect(slide, x, y + Inches(0.45), w, h - Inches(0.45), fill=body_bg)
    bullet_box(slide, body_lines,
               x + Inches(0.15), y + Inches(0.52),
               w - Inches(0.3), h - Inches(0.6),
               size=body_size, color=NAVY)


# ═══════════════════════════════════════════════════════════════════════════════
# Slide 1 — Title
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, W, H, fill=NAVY)
add_rect(sl, 0, Inches(2.8), W, Inches(2.1), fill=RGBColor(0x16, 0x2D, 0x4E))
add_rect(sl, Inches(0.5), Inches(2.75), Inches(0.12), Inches(2.2), fill=GOLD)

add_text(sl, "Nash Analysis", Inches(0.8), Inches(0.6), Inches(11.5), Inches(1.5),
         size=52, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
add_text(sl, "Event Finance Tracker", Inches(0.8), Inches(1.95), Inches(11), Inches(0.7),
         size=26, bold=False, color=GOLD, align=PP_ALIGN.LEFT)
add_text(sl, "Track gross margin, artist costs, and revenue\nacross every show — from import to export.",
         Inches(0.8), Inches(2.9), Inches(11), Inches(1.1),
         size=20, color=RGBColor(0xC5, 0xD5, 0xEA), align=PP_ALIGN.LEFT)
add_text(sl, "Nash Jazz Club  ·  Finance & Operations",
         Inches(0.8), Inches(6.8), Inches(11), Inches(0.45),
         size=13, color=DGRAY, align=PP_ALIGN.LEFT, italic=True)


# ═══════════════════════════════════════════════════════════════════════════════
# Slide 2 — What problem does it solve
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
navy_header(sl, "What Problem Does It Solve?",
            "Consolidating years of event data scattered across Excel files")

add_rect(sl, Inches(0.4), Inches(1.6), Inches(5.9), Inches(5.5), fill=LGRAY)
add_rect(sl, Inches(0.4), Inches(1.6), Inches(5.9), Inches(0.45), fill=RED)
add_text(sl, "Before", Inches(0.55), Inches(1.63), Inches(5.5), Inches(0.38),
         size=16, bold=True, color=WHITE)
bullet_box(sl, [
    "·  Separate Excel workbook per fiscal year & series",
    "·  Each workbook has one tab per month",
    "·  Events are columns; financial rows vary by file",
    "·  No cross-file search or totals",
    "·  Sharing means emailing files — version chaos",
    "·  Artist fee history buried across dozens of tabs",
], Inches(0.55), Inches(2.15), Inches(5.6), Inches(4.7),
   size=15, color=NAVY)

add_rect(sl, Inches(7.0), Inches(1.6), Inches(5.9), Inches(5.5), fill=LGRAY)
add_rect(sl, Inches(7.0), Inches(1.6), Inches(5.9), Inches(0.45), fill=GREEN)
add_text(sl, "After — Nash Analysis", Inches(7.15), Inches(1.63), Inches(5.6), Inches(0.38),
         size=16, bold=True, color=WHITE)
bullet_box(sl, [
    "·  All events in one searchable database",
    "·  Filter by year, series, program, or keyword",
    "·  Sort by revenue, margin, tickets sold, date",
    "·  Artist directory with per-artist totals",
    "·  Password-protected web interface",
    "·  Export back to Excel any time",
], Inches(7.15), Inches(2.15), Inches(5.6), Inches(4.7),
   size=15, color=NAVY)

add_rect(sl, Inches(6.17), Inches(1.6), Inches(0.97), Inches(5.5),
         fill=RGBColor(0xFF, 0xFF, 0xFF))
add_text(sl, "→", Inches(6.17), Inches(3.8), Inches(0.97), Inches(0.9),
         size=40, bold=True, color=NAVY, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# Slide 3 — How data gets in
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
navy_header(sl, "How Data Gets In", "One-time import from existing Excel files")

add_rect(sl, 0, Inches(7.2), W, Inches(0.3), fill=LGRAY)

steps = [
    ("1", "Excel Workbooks", "FY_20xx_*.xlsx files\nOne tab per month\nEvents as columns"),
    ("2", "Convert to CSV", "xlsx_to_csv.py\nDumps every tab to\nper-sheet CSV files"),
    ("3", "Build Database", "build_event_db.py\nPivots & normalizes\ninto SQLite"),
    ("4", "events.db", "Single file\nAll years, all series\n579 events loaded"),
]

arrow_x = [Inches(1.0), Inches(4.3), Inches(7.55), Inches(10.85)]
box_w   = Inches(2.8)
box_h   = Inches(4.2)
top_y   = Inches(1.65)

for idx, (num, title, body) in enumerate(steps):
    bx = arrow_x[idx]
    is_last = idx == len(steps) - 1
    bg = NAVY if is_last else LTBLUE
    title_c = WHITE if is_last else NAVY
    num_c   = GOLD

    add_rect(sl, bx, top_y, box_w, box_h, fill=bg)
    add_text(sl, num, bx + Inches(0.1), top_y + Inches(0.1),
             Inches(0.5), Inches(0.55), size=28, bold=True, color=num_c)
    add_text(sl, title, bx + Inches(0.1), top_y + Inches(0.65),
             box_w - Inches(0.2), Inches(0.55),
             size=17, bold=True, color=title_c)
    add_text(sl, body, bx + Inches(0.15), top_y + Inches(1.25),
             box_w - Inches(0.3), Inches(2.7),
             size=14, color=WHITE if is_last else DGRAY, wrap=True)

    if not is_last:
        ax = bx + box_w + Inches(0.05)
        add_text(sl, "▶", ax, top_y + Inches(1.7),
                 Inches(0.45), Inches(0.6), size=22, bold=True,
                 color=NAVY, align=PP_ALIGN.CENTER)

add_text(sl, "After initial load, new events are entered directly through the web interface.",
         Inches(0.5), Inches(6.3), Inches(12), Inches(0.5),
         size=14, color=DGRAY, italic=True, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# Slide 4 — Event List (main screen)
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
navy_header(sl, "Event List — Main Screen",
            "Search, filter, and sort every show at a glance")

add_rect(sl, Inches(0.4), Inches(1.55), Inches(12.5), Inches(5.65), fill=LGRAY)

# filter bar mock
add_rect(sl, Inches(0.5), Inches(1.7), Inches(12.3), Inches(0.55),
         fill=RGBColor(0xE5, 0xE7, 0xEB))
for label, xp in [("Search events…", 0.6), ("Series ▾", 4.2),
                   ("Program ▾", 5.9), ("FY ▾", 7.6)]:
    add_text(sl, label, Inches(xp), Inches(1.75), Inches(1.5), Inches(0.42),
             size=11, color=DGRAY)

# table header
cols   = ["Date", "Name", "Series", "Revenue", "Expense", "Margin", "Tickets", ""]
widths = [1.1,    3.5,    1.5,      1.2,        1.2,       1.2,      1.0,       0.5]
starts = [0.5]
for w in widths[:-1]:
    starts.append(starts[-1] + w)

add_rect(sl, Inches(0.5), Inches(2.35), Inches(12.3), Inches(0.38), fill=NAVY)
for c, label in enumerate(cols):
    add_text(sl, label, Inches(starts[c] + 0.05), Inches(2.38),
             Inches(widths[c]), Inches(0.32),
             size=11, bold=True, color=WHITE)

rows = [
    ("Apr 26, 2026", "Kevin Hays Trio", "Performance", "$7,820", "$2,340", "$5,480", "148"),
    ("Apr 10, 2026", "Endea Owens (Early)", "Performance", "$5,120", "$3,800", "$1,320", "96"),
    ("Apr 10, 2026", "Endea Owens (Late)", "Performance", "$4,990", "$3,800", "$1,190", "91"),
    ("Mar 28, 2026", "First Friday March", "First Friday", "$3,240", "$1,100", "$2,140", "112"),
    ("Mar 14, 2026", "Education Show", "Education", "$0", "$950", "($950)", "44"),
]

row_bgs = [WHITE, LGRAY]
for r, row_data in enumerate(rows):
    ry = Inches(2.73 + r * 0.52)
    add_rect(sl, Inches(0.5), ry, Inches(12.3), Inches(0.5),
             fill=row_bgs[r % 2])
    vals = list(row_data) + ["✕"]
    for c, val in enumerate(vals):
        color = NAVY
        if c == 5:
            color = GREEN if not val.startswith("(") else RED
        add_text(sl, val, Inches(starts[c] + 0.05), ry + Inches(0.1),
                 Inches(widths[c]), Inches(0.35),
                 size=12, color=color,
                 bold=(c == 5))

add_text(sl, "Total (579 events)",
         Inches(0.55), Inches(5.6), Inches(3.5), Inches(0.4),
         size=12, bold=True, color=NAVY)
add_text(sl, "$4,128,440    $2,943,210    $1,185,230",
         Inches(4.0), Inches(5.6), Inches(5.0), Inches(0.4),
         size=12, bold=True, color=NAVY)

bullet_box(sl, [
    "·  Filter simultaneously by keyword, series, program, and fiscal year",
    "·  Click any column header to sort  ·  Totals row always reflects current filter",
], Inches(0.5), Inches(6.3), Inches(12.3), Inches(0.7),
   size=13, color=DGRAY)


# ═══════════════════════════════════════════════════════════════════════════════
# Slide 5 — Event Detail & Edit
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
navy_header(sl, "Event Detail & Edit",
            "Full financial breakdown and in-browser editing")

# left panel — detail
add_rect(sl, Inches(0.4), Inches(1.55), Inches(6.0), Inches(5.65), fill=LGRAY)
add_rect(sl, Inches(0.4), Inches(1.55), Inches(6.0), Inches(0.42), fill=NAVY)
add_text(sl, "Kevin Hays Trio — Apr 26, 2026", Inches(0.55), Inches(1.58),
         Inches(5.7), Inches(0.36), size=13, bold=True, color=WHITE)

fields = [
    ("Series",   "Performance"),
    ("Program",  "Performance"),
    ("FY",       "2026"),
    ("Artists",  "Kevin Hays"),
    ("", ""),
    ("REVENUE", ""),
    ("Ticket Sales", "$7,820.00"),
    ("F&B Revenue",  "$0.00"),
    ("Donations",    "$0.00"),
    ("Total Revenue","$7,820.00"),
    ("", ""),
    ("EXPENSE", ""),
    ("Performance Fee", "$2,100.00"),
    ("Tech Support",    "$240.00"),
    ("Total Expense",   "$2,340.00"),
    ("", ""),
    ("Gross Margin",    "$5,480.00"),
]
for fi, (k, v) in enumerate(fields):
    fy2 = Inches(2.08 + fi * 0.25)
    if k in ("REVENUE", "EXPENSE"):
        add_text(sl, k, Inches(0.55), fy2, Inches(5.5), Inches(0.23),
                 size=10, bold=True, color=DGRAY)
        continue
    if not k:
        continue
    bold_row = k in ("Total Revenue", "Total Expense", "Gross Margin")
    col_v = GREEN if k == "Gross Margin" else NAVY
    add_text(sl, k, Inches(0.55), fy2, Inches(3.2), Inches(0.23),
             size=11, bold=bold_row, color=DGRAY if not bold_row else NAVY)
    add_text(sl, v, Inches(3.8), fy2, Inches(2.4), Inches(0.23),
             size=11, bold=bold_row, color=col_v, align=PP_ALIGN.RIGHT)

# right panel — edit form
add_rect(sl, Inches(6.9), Inches(1.55), Inches(6.0), Inches(5.65), fill=LGRAY)
add_rect(sl, Inches(6.9), Inches(1.55), Inches(6.0), Inches(0.42), fill=GOLD)
add_text(sl, "Edit Event", Inches(7.05), Inches(1.58),
         Inches(5.7), Inches(0.36), size=13, bold=True, color=NAVY)

form_fields = [
    ("Date", "2026-04-26"),
    ("Name", "Kevin Hays Trio"),
    ("Series", "Performance"),
    ("Artists", "Kevin Hays"),
    ("Ticket Sales", "7820"),
    ("Performance Fee", "2100"),
    ("Tech Support", "240"),
    ("Tickets Sold", "148"),
    ("Notes", ""),
]
for fi2, (label, val) in enumerate(form_fields):
    gy = Inches(2.12 + fi2 * 0.5)
    add_text(sl, label, Inches(7.05), gy, Inches(2.2), Inches(0.26),
             size=11, color=DGRAY)
    add_rect(sl, Inches(9.3), gy - Inches(0.03), Inches(3.3), Inches(0.32),
             fill=WHITE)
    add_text(sl, val, Inches(9.38), gy, Inches(3.1), Inches(0.28),
             size=11, color=NAVY)

add_rect(sl, Inches(9.3), Inches(6.55), Inches(1.5), Inches(0.35), fill=NAVY)
add_text(sl, "Save", Inches(9.3), Inches(6.57), Inches(1.5), Inches(0.3),
         size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_rect(sl, Inches(11.0), Inches(6.55), Inches(1.6), Inches(0.35), fill=LGRAY)
add_text(sl, "Duplicate", Inches(11.0), Inches(6.57), Inches(1.6), Inches(0.3),
         size=12, bold=False, color=NAVY, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# Slide 6 — Artist Directory
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
navy_header(sl, "Artist Directory",
            "Per-artist revenue totals across all shows and years")

add_rect(sl, Inches(0.4), Inches(1.55), Inches(7.8), Inches(5.65), fill=LGRAY)
add_rect(sl, Inches(0.4), Inches(1.55), Inches(7.8), Inches(0.42), fill=NAVY)
for label, xp in [("Artist", 0.55), ("Shows", 4.5), ("Total Revenue", 5.7)]:
    add_text(sl, label, Inches(xp), Inches(1.6), Inches(1.8), Inches(0.33),
             size=12, bold=True, color=WHITE)

artist_rows = [
    ("Kevin Hays", "14", "$48,320"),
    ("Endea Owens", "9", "$38,740"),
    ("Moe Flavour", "11", "$31,200"),
    ("Francine Reed", "8", "$27,880"),
    ("Jocelyn Gould Quartet", "6", "$22,140"),
    ("Sonrisa", "7", "$18,990"),
    ("Dave Potter", "12", "$17,650"),
    ("Various Artists", "22", "$14,300"),
]
for r, (name, shows, rev) in enumerate(artist_rows):
    ry = Inches(2.08 + r * 0.45)
    add_rect(sl, Inches(0.4), ry, Inches(7.8), Inches(0.43),
             fill=WHITE if r % 2 == 0 else LGRAY)
    add_text(sl, name, Inches(0.55), ry + Inches(0.07),
             Inches(3.8), Inches(0.32), size=12, color=NAVY)
    add_text(sl, shows, Inches(4.5), ry + Inches(0.07),
             Inches(1.0), Inches(0.32), size=12, color=DGRAY,
             align=PP_ALIGN.CENTER)
    add_text(sl, rev, Inches(5.7), ry + Inches(0.07),
             Inches(2.2), Inches(0.32), size=12, bold=True, color=GREEN,
             align=PP_ALIGN.RIGHT)

# right — artist detail card
add_rect(sl, Inches(8.7), Inches(1.55), Inches(4.6), Inches(5.65), fill=LTBLUE)
add_rect(sl, Inches(8.7), Inches(1.55), Inches(4.6), Inches(0.42), fill=GOLD)
add_text(sl, "Kevin Hays — Artist Detail", Inches(8.85), Inches(1.58),
         Inches(4.3), Inches(0.36), size=12, bold=True, color=NAVY)

kh_rows = [
    ("Apr 26, 2026", "Trio", "$5,480"),
    ("Jan 18, 2026", "Solo", "$3,200"),
    ("Sep 12, 2025", "Trio", "$4,910"),
    ("Jun 7, 2025",  "Trio", "$4,750"),
    ("Mar 15, 2025", "Quartet", "$3,980"),
]
add_rect(sl, Inches(8.7), Inches(2.05), Inches(4.6), Inches(0.35), fill=NAVY)
for lbl, xp2 in [("Date", 8.85), ("Show", 10.1), ("Margin", 11.6)]:
    add_text(sl, lbl, Inches(xp2), Inches(2.08), Inches(1.4), Inches(0.28),
             size=11, bold=True, color=WHITE)

for r2, (date, show, margin) in enumerate(kh_rows):
    ry2 = Inches(2.48 + r2 * 0.44)
    add_rect(sl, Inches(8.7), ry2, Inches(4.6), Inches(0.42),
             fill=WHITE if r2 % 2 == 0 else LGRAY)
    add_text(sl, date, Inches(8.85), ry2 + Inches(0.08),
             Inches(1.2), Inches(0.3), size=11, color=NAVY)
    add_text(sl, show, Inches(10.1), ry2 + Inches(0.08),
             Inches(1.4), Inches(0.3), size=11, color=DGRAY)
    add_text(sl, margin, Inches(11.6), ry2 + Inches(0.08),
             Inches(1.5), Inches(0.3), size=11, bold=True, color=GREEN,
             align=PP_ALIGN.RIGHT)

bullet_box(sl, ["Total: 14 shows  ·  $48,320 revenue  ·  $29,840 margin"],
           Inches(8.85), Inches(4.73), Inches(4.3), Inches(0.4),
           size=12, color=NAVY)

add_text(sl, "Revenue reflects each artist's share % of the event.",
         Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.5),
         size=13, color=DGRAY, italic=True, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# Slide 7 — Excel Export
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
navy_header(sl, "Excel Export — Getting Data Back Out",
            "One command produces a fully-formatted workbook")

# left: command + description
add_rect(sl, Inches(0.4), Inches(1.6), Inches(5.8), Inches(5.6), fill=LGRAY)

add_rect(sl, Inches(0.5), Inches(1.75), Inches(5.6), Inches(0.95),
         fill=RGBColor(0x1E, 0x1E, 0x1E))
add_text(sl, "python export_to_excel.py",
         Inches(0.6), Inches(1.82), Inches(5.4), Inches(0.4),
         size=14, bold=True, color=RGBColor(0x9C, 0xD5, 0xFF))
add_text(sl, "→  events_export.xlsx",
         Inches(0.6), Inches(2.2), Inches(5.4), Inches(0.35),
         size=12, color=RGBColor(0x6A, 0x99, 0x55))

bullet_box(sl, [
    "·  One row per event — all 579 events",
    "·  Columns: Date, Name, Series, Program,",
    "   FY, Artists, all revenue & expense fields",
    "·  Total Revenue and Total Expense are live",
    "   Excel SUM() formulas — edit a cell and",
    "   totals recalculate automatically",
    "·  Gross Margin = Total Revenue − Total Expense",
    "·  Revenue cells shaded green, expense red",
    "·  Bold margin column; negative = red text",
    "·  Ticket counts and Notes included",
], Inches(0.55), Inches(2.85), Inches(5.6), Inches(4.1),
   size=13, color=NAVY)

# right: column group visual
add_rect(sl, Inches(6.6), Inches(1.6), Inches(6.3), Inches(5.6), fill=LGRAY)
add_text(sl, "Column layout", Inches(6.75), Inches(1.68), Inches(6.0), Inches(0.4),
         size=14, bold=True, color=NAVY)

groups = [
    (NAVY,  "Event Info",     "Date · Name · Series · Program · FY · Artists"),
    (GREEN, "Revenue (×8)",   "Ticket Sales · F&B · Merch · CDs · Donations\nGrants · Venue Rental · Other  →  Total Revenue"),
    (RED,   "Expense (×9)",   "Perf. Fee · Tech · Security · Travel · Equipment\nAdvertising · Misc · Merchant Fees · F&B Exp  →  Total Expense"),
    (NAVY,  "Margin",         "= Total Revenue − Total Expense   (live formula)"),
    (DGRAY, "Counts & Notes", "Tickets Sold · Full Price · Discount · Comp · Notes"),
]
for gi, (color, title, desc) in enumerate(groups):
    gy = Inches(2.18 + gi * 0.95)
    add_rect(sl, Inches(6.7), gy, Inches(0.18), Inches(0.72), fill=color)
    add_text(sl, title, Inches(7.0), gy, Inches(5.6), Inches(0.32),
             size=13, bold=True, color=NAVY)
    add_text(sl, desc, Inches(7.0), gy + Inches(0.32), Inches(5.7), Inches(0.5),
             size=11, color=DGRAY, wrap=True)

add_text(sl, "--db and --out flags let you target a different database or filename.",
         Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.5),
         size=13, color=DGRAY, italic=True, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# Slide 8 — How to Run It
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
navy_header(sl, "How to Run It", "Local and server deployment options")

# local card
add_rect(sl, Inches(0.4), Inches(1.6), Inches(5.9), Inches(5.55), fill=LGRAY)
add_rect(sl, Inches(0.4), Inches(1.6), Inches(5.9), Inches(0.44), fill=NAVY)
add_text(sl, "Local (laptop)", Inches(0.55), Inches(1.63),
         Inches(5.6), Inches(0.38), size=16, bold=True, color=WHITE)

local_cmds = [
    ("# activate environment", DGRAY),
    ("source .venv/bin/activate", RGBColor(0x9C, 0xD5, 0xFF)),
    ("", WHITE),
    ("# start the app", DGRAY),
    ("python app.py", RGBColor(0x9C, 0xD5, 0xFF)),
    ("", WHITE),
    ("# custom port or database", DGRAY),
    ("python app.py --port 8000", RGBColor(0x9C, 0xD5, 0xFF)),
    ("python app.py --db backup.db", RGBColor(0x9C, 0xD5, 0xFF)),
]
add_rect(sl, Inches(0.5), Inches(2.16), Inches(5.7), Inches(2.6),
         fill=RGBColor(0x1E, 0x1E, 0x1E))
for ci, (cmd, color) in enumerate(local_cmds):
    add_text(sl, cmd, Inches(0.62), Inches(2.24 + ci * 0.26),
             Inches(5.5), Inches(0.25), size=11, color=color)

bullet_box(sl, [
    "·  Opens at http://127.0.0.1:5050",
    "·  Login with username + password",
    "·  No internet connection required",
    "·  Database is a single local file (events.db)",
], Inches(0.55), Inches(4.9), Inches(5.6), Inches(2.0), size=13, color=NAVY)

# server card
add_rect(sl, Inches(7.0), Inches(1.6), Inches(5.9), Inches(5.55), fill=LGRAY)
add_rect(sl, Inches(7.0), Inches(1.6), Inches(5.9), Inches(0.44), fill=GOLD)
add_text(sl, "Server (shared access)", Inches(7.15), Inches(1.63),
         Inches(5.6), Inches(0.38), size=16, bold=True, color=NAVY)

server_cmds = [
    ("# one-command deploy via Ansible", DGRAY),
    ("./deploy.sh", RGBColor(0x9C, 0xD5, 0xFF)),
    ("", WHITE),
    ("# or with Docker Compose directly", DGRAY),
    ("docker-compose up -d", RGBColor(0x9C, 0xD5, 0xFF)),
]
add_rect(sl, Inches(7.1), Inches(2.16), Inches(5.7), Inches(1.65),
         fill=RGBColor(0x1E, 0x1E, 0x1E))
for ci2, (cmd, color) in enumerate(server_cmds):
    add_text(sl, cmd, Inches(7.22), Inches(2.24 + ci2 * 0.26),
             Inches(5.5), Inches(0.25), size=11, color=color)

bullet_box(sl, [
    "·  Runs in Docker — no manual install on server",
    "·  Accessible from any browser on the network",
    "·  Password protection built in",
    "·  Ansible playbook handles first-time setup",
    "   and re-deploys with a single command",
    "·  Secrets stored in encrypted Ansible Vault",
], Inches(7.15), Inches(3.95), Inches(5.6), Inches(3.0), size=13, color=NAVY)


# ═══════════════════════════════════════════════════════════════════════════════
# Slide 9 — Summary
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, W, H, fill=NAVY)
add_rect(sl, 0, Inches(3.35), W, Inches(0.07), fill=GOLD)

add_text(sl, "Nash Analysis at a Glance",
         Inches(0.6), Inches(0.35), Inches(12), Inches(0.9),
         size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

summary_items = [
    ("579", "events loaded\nacross all years"),
    ("20+", "financial fields\nper event"),
    ("1", "file to back up\n(events.db)"),
    ("1", "command to export\nback to Excel"),
]
for si, (big, small) in enumerate(summary_items):
    sx = Inches(0.6 + si * 3.2)
    add_rect(sl, sx, Inches(1.4), Inches(2.9), Inches(1.65),
             fill=RGBColor(0x16, 0x2D, 0x4E))
    add_text(sl, big, sx + Inches(0.1), Inches(1.48),
             Inches(1.0), Inches(0.9), size=42, bold=True, color=GOLD)
    add_text(sl, small, sx + Inches(1.1), Inches(1.6),
             Inches(1.7), Inches(0.8), size=13, color=RGBColor(0xC5, 0xD5, 0xEA),
             wrap=True)

caps = [
    ("Import", "Excel → CSV → SQLite\none-time bulk load"),
    ("Search & Filter", "By year, series, program,\nor keyword"),
    ("Edit In-Browser", "Add, edit, duplicate,\nor delete any event"),
    ("Artist Tracking", "Per-artist totals\nacross all shows"),
    ("Export", "Back to .xlsx\nwith live formulas"),
    ("Deploy", "Local or server\nvia Docker + Ansible"),
]
for ci3, (title, desc) in enumerate(caps):
    cx = Inches(0.5 + (ci3 % 3) * 4.25)
    cy = Inches(3.6 + (ci3 // 3) * 1.55)
    add_rect(sl, cx, cy, Inches(3.9), Inches(1.35),
             fill=RGBColor(0x16, 0x2D, 0x4E))
    add_rect(sl, cx, cy, Inches(0.12), Inches(1.35), fill=GOLD)
    add_text(sl, title, cx + Inches(0.22), cy + Inches(0.1),
             Inches(3.5), Inches(0.38), size=14, bold=True, color=WHITE)
    add_text(sl, desc, cx + Inches(0.22), cy + Inches(0.5),
             Inches(3.5), Inches(0.75), size=12,
             color=RGBColor(0xC5, 0xD5, 0xEA), wrap=True)


prs.save(out_path)
print(f"Saved → {out_path}")
