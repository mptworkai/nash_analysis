#!/usr/bin/env python3
"""
Generate a pitch deck for the Nash Analysis web front end.
Audience: stakeholders / staff who have not seen the tool before.

Usage:
  python make_pitch.py                  # writes nash_pitch_deck.pptx
  python make_pitch.py --out FILE.pptx
  python make_pitch.py --help
"""

import sys

if "-h" in sys.argv or "--help" in sys.argv:
    print(__doc__)
    sys.exit(0)

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
except ImportError:
    print("Missing dependency. Run: pip install python-pptx", file=sys.stderr)
    sys.exit(1)

from pathlib import Path

out_path = Path("nash_pitch_deck.pptx")
args = sys.argv[1:]
i = 0
while i < len(args):
    if args[i] == "--out" and i + 1 < len(args):
        out_path = Path(args[i + 1]); i += 2
    else:
        i += 1

# ── palette ───────────────────────────────────────────────────────────────────
NAVY    = RGBColor(0x1E, 0x3A, 0x5F)
NAVY2   = RGBColor(0x16, 0x2D, 0x4E)
GOLD    = RGBColor(0xC8, 0x9B, 0x2E)
GOLD2   = RGBColor(0xF5, 0xD8, 0x7A)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY   = RGBColor(0xF3, 0xF4, 0xF6)
MGRAY   = RGBColor(0xE5, 0xE7, 0xEB)
DGRAY   = RGBColor(0x6B, 0x72, 0x80)
GREEN   = RGBColor(0x06, 0x5F, 0x46)
LGREEN  = RGBColor(0xD1, 0xFA, 0xE5)
RED     = RGBColor(0x7F, 0x1D, 0x1D)
LRED    = RGBColor(0xFE, 0xE2, 0xE2)
AMBER   = RGBColor(0x92, 0x40, 0x0E)
LAMBER  = RGBColor(0xFF, 0xED, 0xD5)

W = Inches(13.333)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]


# ── helpers ───────────────────────────────────────────────────────────────────

def rect(sl, x, y, w, h, fill=None):
    s = sl.shapes.add_shape(1, x, y, w, h)
    if fill:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    s.line.fill.background()
    return s


def txt(sl, text, x, y, w, h, size=18, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, italic=False, wrap=True):
    tb = sl.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = wrap
    p  = tf.paragraphs[0]; p.alignment = align
    r  = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = color
    return tb


def multiline(sl, lines, x, y, w, h, size=15, color=NAVY, bold=False,
              line_space=None):
    tb = sl.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    first = True
    for line in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        if line_space:
            p.space_before = Pt(line_space)
        r = p.add_run(); r.text = line
        r.font.size = Pt(size); r.font.color.rgb = color
        r.font.bold = bold


def divider(sl, y, color=GOLD):
    rect(sl, 0, y, W, Inches(0.05), fill=color)


def stat_block(sl, x, y, w, h, number, label, bg=NAVY2, num_color=GOLD,
               lbl_color=WHITE):
    rect(sl, x, y, w, h, fill=bg)
    txt(sl, number, x, y + Inches(0.08), w, Inches(0.75),
        size=40, bold=True, color=num_color, align=PP_ALIGN.CENTER)
    txt(sl, label,  x, y + Inches(0.75), w, Inches(0.55),
        size=13, color=lbl_color, align=PP_ALIGN.CENTER, wrap=True)


def pain_card(sl, x, y, w, h, emoji, headline, body):
    rect(sl, x, y, w, h, fill=LRED)
    rect(sl, x, y, w, Inches(0.06), fill=RED)
    txt(sl, emoji,    x, y + Inches(0.12), w, Inches(0.55),
        size=28, align=PP_ALIGN.CENTER, color=RED)
    txt(sl, headline, x, y + Inches(0.65), w, Inches(0.45),
        size=14, bold=True, color=RED, align=PP_ALIGN.CENTER)
    txt(sl, body,     x + Inches(0.15), y + Inches(1.15),
        w - Inches(0.3), h - Inches(1.25),
        size=12, color=AMBER, wrap=True)


def feature_card(sl, x, y, w, h, icon, title, body, accent=NAVY):
    rect(sl, x, y, w, h, fill=LGRAY)
    rect(sl, x, y, Inches(0.1), h, fill=accent)
    txt(sl, icon,  x + Inches(0.25), y + Inches(0.12),
        Inches(0.7), Inches(0.55), size=26, color=accent)
    txt(sl, title, x + Inches(0.25), y + Inches(0.62),
        w - Inches(0.35), Inches(0.38),
        size=14, bold=True, color=NAVY)
    txt(sl, body,  x + Inches(0.25), y + Inches(1.0),
        w - Inches(0.35), h - Inches(1.1),
        size=12, color=DGRAY, wrap=True)


# ═══════════════════════════════════════════════════════════════════════════════
# Slide 1 — Title / Hook
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
rect(sl, 0, 0, W, H, fill=NAVY)

# gold accent bar left
rect(sl, 0, 0, Inches(0.45), H, fill=GOLD)

# big hook line
txt(sl, "Every show.\nEvery dollar.\nOne place.",
    Inches(0.9), Inches(0.6), Inches(7.5), Inches(3.5),
    size=54, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

txt(sl, "Introducing Nash Analysis —\na purpose-built finance tracker for Nash Jazz Club events.",
    Inches(0.9), Inches(4.3), Inches(7.5), Inches(1.1),
    size=20, color=GOLD2, align=PP_ALIGN.LEFT)

txt(sl, "Gross margin  ·  Artist costs  ·  Revenue trends  ·  Instant search",
    Inches(0.9), Inches(5.55), Inches(7.5), Inches(0.5),
    size=14, italic=True, color=DGRAY, align=PP_ALIGN.LEFT)

# stat column right
for si, (num, lbl) in enumerate([
    ("579",  "events tracked"),
    ("FY22–", "through today"),
    ("20+",  "financial fields\nper event"),
]):
    stat_block(sl, Inches(9.8), Inches(0.9 + si * 2.1),
               Inches(3.1), Inches(1.8), num, lbl)

txt(sl, "Nash Jazz Club  ·  Finance & Operations",
    Inches(0.9), Inches(7.0), Inches(8.0), Inches(0.38),
    size=12, italic=True, color=DGRAY)


# ═══════════════════════════════════════════════════════════════════════════════
# Slide 2 — The Problem
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
rect(sl, 0, 0, W, Inches(1.25), fill=NAVY)
txt(sl, "The Problem", Inches(0.5), Inches(0.18),
    Inches(10), Inches(0.65), size=34, bold=True, color=WHITE)
txt(sl, "Managing event finances across dozens of Excel files is slow, fragile, and error-prone.",
    Inches(0.5), Inches(0.78), Inches(12.5), Inches(0.4),
    size=16, color=GOLD2)
divider(sl, Inches(1.25))

cards = [
    ("📁", "Scattered Files",
     "One workbook per fiscal year and series.\nEach workbook has one tab per month.\nFinding a single show means opening\nthe right file, the right tab."),
    ("🔀", "No Consistency",
     "Column layouts shift between years.\nMetric names vary (\"Performers Fee\"\nvs. \"Performance Fee\"). Formulas\ndrift from file to file."),
    ("🔍", "Zero Cross-File Search",
     "Can't answer \"What did we pay this\nartist across all years?\" without\nopening every file manually\nand summing by hand."),
    ("📧", "Version Chaos",
     "Sharing means emailing attachments.\nWhich version is current? Who changed\nrow 14? There's no way to know\nwithout asking."),
    ("🧮", "Manual Totals",
     "Annual and multi-series summaries\nrequire copying numbers between files.\nOne wrong cell and the margin\nnumbers are wrong all year."),
    ("⏱️", "Slow to Answer Questions",
     "\"How was First Friday vs. Special\nEvents this year?\" takes 20 minutes\nof file-hopping instead of a\nsingle search."),
]

cw = Inches(3.9)
ch = Inches(3.55)
for ci, (emoji, headline, body) in enumerate(cards):
    col = ci % 3; row = ci // 3
    cx = Inches(0.4 + col * 4.3)
    cy = Inches(1.5 + row * 3.75)
    if cy + ch > H:
        cy = H - ch - Inches(0.1)
    pain_card(sl, cx, cy, cw, ch, emoji, headline, body)


# ═══════════════════════════════════════════════════════════════════════════════
# Slide 3 — The Cost
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
rect(sl, 0, 0, W, Inches(1.25), fill=NAVY)
txt(sl, "What It's Costing Us", Inches(0.5), Inches(0.18),
    Inches(10), Inches(0.65), size=34, bold=True, color=WHITE)
txt(sl, "Invisible overhead that compounds every month.",
    Inches(0.5), Inches(0.78), Inches(12.5), Inches(0.4),
    size=16, color=GOLD2)
divider(sl, Inches(1.25))

costs = [
    ("Time", "30–60 min",
     "to answer a question that should take 30 seconds —\n"
     "\"What did we gross on jazz performances in FY2025?\""),
    ("Errors", "Silent mistakes",
     "copy-paste errors and broken formulas that go\n"
     "undetected until the annual review — or longer"),
    ("Decisions", "Delayed insights",
     "pricing, booking, and program decisions made\n"
     "without timely margin data because pulling it\n"
     "takes too long"),
    ("Knowledge", "People-dependent",
     "only the person who built the spreadsheets knows\n"
     "where everything lives — that's a single point of failure"),
]

for ki, (label, impact, desc) in enumerate(costs):
    ky = Inches(1.55 + ki * 1.42)
    rect(sl, Inches(0.4), ky, Inches(12.5), Inches(1.28), fill=LRED)
    rect(sl, Inches(0.4), ky, Inches(0.08), Inches(1.28), fill=RED)
    txt(sl, label,  Inches(0.65), ky + Inches(0.12),
        Inches(1.8), Inches(0.4), size=14, bold=True, color=RED)
    txt(sl, impact, Inches(2.5),  ky + Inches(0.1),
        Inches(2.8), Inches(0.5), size=22, bold=True, color=AMBER)
    txt(sl, desc,   Inches(5.4),  ky + Inches(0.18),
        Inches(7.3), Inches(0.95), size=14, color=AMBER, wrap=True)

txt(sl, "None of this is acceptable when the data already exists — it just isn't organized.",
    Inches(0.4), Inches(7.1), Inches(12.5), Inches(0.35),
    size=14, italic=True, color=DGRAY, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# Slide 4 — The Solution
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
rect(sl, 0, 0, W, H, fill=NAVY)
divider(sl, Inches(1.3), color=GOLD)

txt(sl, "Introducing Nash Analysis",
    Inches(0.6), Inches(0.15), Inches(12), Inches(0.85),
    size=38, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txt(sl, "A secure, browser-based finance tracker purpose-built for Nash events.",
    Inches(0.6), Inches(0.95), Inches(12), Inches(0.4),
    size=18, color=GOLD2, align=PP_ALIGN.CENTER)

points = [
    ("All your data, one place",
     "Every show from FY2022 through today loaded and searchable in seconds."),
    ("Built around how you work",
     "Same financial categories as your existing Excel files — no relearning."),
    ("Answers questions instantly",
     "Filter by year, series, or program. Sort by margin, revenue, or tickets. Done."),
    ("Nothing to install",
     "Opens in any browser. Works on a laptop at the office or a phone backstage."),
    ("Your data stays yours",
     "Runs on Nash's own server. Password-protected. No third-party cloud required."),
    ("Excel is still there",
     "Export the full database back to a formatted .xlsx file any time — one command."),
]

for pi, (head, body) in enumerate(points):
    col = pi % 2; row = pi // 2
    px = Inches(0.5 + col * 6.5)
    py = Inches(1.55 + row * 1.8)
    rect(sl, px, py, Inches(6.1), Inches(1.65), fill=NAVY2)
    rect(sl, px, py, Inches(0.1), Inches(1.65), fill=GOLD)
    txt(sl, head, px + Inches(0.25), py + Inches(0.15),
        Inches(5.7), Inches(0.42), size=16, bold=True, color=WHITE)
    txt(sl, body, px + Inches(0.25), py + Inches(0.58),
        Inches(5.7), Inches(0.98), size=13, color=RGBColor(0xC5, 0xD5, 0xEA),
        wrap=True)


# ═══════════════════════════════════════════════════════════════════════════════
# Slide 5 — Feature: Event List
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
rect(sl, 0, 0, W, Inches(1.25), fill=NAVY)
txt(sl, "Feature 1 of 4  —  Event List", Inches(0.5), Inches(0.12),
    Inches(10), Inches(0.5), size=14, color=GOLD)
txt(sl, "Find any show in seconds", Inches(0.5), Inches(0.52),
    Inches(12), Inches(0.65), size=32, bold=True, color=WHITE)
divider(sl, Inches(1.25))

# screen mock
rect(sl, Inches(0.4), Inches(1.4), Inches(8.8), Inches(5.85), fill=MGRAY)

# filter bar
rect(sl, Inches(0.5), Inches(1.55), Inches(8.6), Inches(0.5), fill=WHITE)
for lbl, xp in [("Search events…", 0.62), ("Series ▾", 3.8),
                ("Program ▾", 5.3), ("FY ▾", 6.7), ("+ New Event", 7.5)]:
    c = NAVY if lbl == "+ New Event" else DGRAY
    txt(sl, lbl, Inches(xp), Inches(1.6), Inches(1.4), Inches(0.38),
        size=11, color=c)

# table header
cols   = ["Date",  "Event Name",              "Series",      "Revenue", "Expense", "Margin",  "Tickets"]
widths = [0.95,    3.0,                        1.3,           1.1,       1.1,       1.1,       0.85]
starts = [0.5]
for ww in widths[:-1]:
    starts.append(starts[-1] + ww)

rect(sl, Inches(0.5), Inches(2.13), Inches(8.6), Inches(0.38), fill=NAVY)
for ci, lbl in enumerate(cols):
    txt(sl, lbl, Inches(starts[ci] + 0.05), Inches(2.17),
        Inches(widths[ci]), Inches(0.3), size=10, bold=True, color=WHITE)

rows = [
    ("Apr 26 '26", "Kevin Hays Trio",             "Performance",  "$7,820", "$2,340", "+$5,480", "148"),
    ("Apr 10 '26", "Endea Owens (Early)",          "Performance",  "$5,120", "$3,800", "+$1,320",  "96"),
    ("Apr 10 '26", "Endea Owens (Late)",           "Performance",  "$4,990", "$3,800", "+$1,190",  "91"),
    ("Mar 28 '26", "First Friday March",           "First Friday", "$3,240", "$1,100", "+$2,140", "112"),
    ("Mar 14 '26", "Education Matinee",            "Education",    "$0",     "$950",   "($950)",   "44"),
    ("Feb 21 '26", "Moe Flavour Quartet",          "Performance",  "$6,110", "$2,800", "+$3,310", "134"),
    ("Jan 18 '26", "First Friday January",         "First Friday", "$2,980", "$900",   "+$2,080", "108"),
    ("Jan  4 '26", "New Year Kickoff Special",     "Special Event","$9,400", "$4,200", "+$5,200", "210"),
]
for ri, row in enumerate(rows):
    ry = Inches(2.56 + ri * 0.44)
    rect(sl, Inches(0.5), ry, Inches(8.6), Inches(0.42),
         fill=WHITE if ri % 2 == 0 else LGRAY)
    for ci2, val in enumerate(row):
        c = NAVY
        if ci2 == 5:
            c = GREEN if val.startswith("+") else RED
        txt(sl, val, Inches(starts[ci2] + 0.05), ry + Inches(0.09),
            Inches(widths[ci2]), Inches(0.3),
            size=11, color=c, bold=(ci2 == 5))

# totals
rect(sl, Inches(0.5), Inches(7.05), Inches(8.6), Inches(0.38), fill=NAVY2)
txt(sl, "Total — 579 events", Inches(0.6), Inches(7.08),
    Inches(3.5), Inches(0.3), size=11, bold=True, color=WHITE)
txt(sl, "$4,128,440   $2,943,210   +$1,185,230",
    Inches(4.1), Inches(7.08), Inches(4.8), Inches(0.3),
    size=11, bold=True, color=GOLD2)

# call-outs right
callouts = [
    ("Filter simultaneously", "Search by keyword AND series AND\nprogram AND fiscal year at once."),
    ("Sort any column",       "Click a header — sort by margin,\nrevenue, tickets, or date."),
    ("Totals always match",   "The summary row at the bottom\nreflects exactly the current filter."),
    ("Instant delete",        "Remove an event from the list\nwith one confirmed click."),
]
for ci3, (head, body) in enumerate(callouts):
    cy2 = Inches(1.4 + ci3 * 1.52)
    rect(sl, Inches(9.55), cy2, Inches(3.45), Inches(1.38), fill=LGREEN)
    rect(sl, Inches(9.55), cy2, Inches(0.1), Inches(1.38), fill=GREEN)
    txt(sl, head, Inches(9.75), cy2 + Inches(0.1),
        Inches(3.1), Inches(0.36), size=13, bold=True, color=GREEN)
    txt(sl, body, Inches(9.75), cy2 + Inches(0.48),
        Inches(3.1), Inches(0.78), size=12, color=NAVY, wrap=True)


# ═══════════════════════════════════════════════════════════════════════════════
# Slide 6 — Feature: Event Detail & Edit
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
rect(sl, 0, 0, W, Inches(1.25), fill=NAVY)
txt(sl, "Feature 2 of 4  —  Event Detail & Edit", Inches(0.5), Inches(0.12),
    Inches(10), Inches(0.5), size=14, color=GOLD)
txt(sl, "See everything. Change anything.", Inches(0.5), Inches(0.52),
    Inches(12), Inches(0.65), size=32, bold=True, color=WHITE)
divider(sl, Inches(1.25))

# detail panel
rect(sl, Inches(0.4), Inches(1.4), Inches(4.5), Inches(5.9), fill=LGRAY)
rect(sl, Inches(0.4), Inches(1.4), Inches(4.5), Inches(0.45), fill=NAVY)
txt(sl, "Kevin Hays Trio  —  Apr 26, 2026",
    Inches(0.55), Inches(1.44), Inches(4.2), Inches(0.38),
    size=13, bold=True, color=WHITE)

detail_rows = [
    ("Series",          "Performance",  False, DGRAY),
    ("Program",         "Performance",  False, DGRAY),
    ("Fiscal Year",     "2026",         False, DGRAY),
    ("Artists",         "Kevin Hays",   False, DGRAY),
    ("",                "",             False, DGRAY),
    ("REVENUE",         "",             True,  DGRAY),
    ("Ticket Sales",    "$7,820",       False, DGRAY),
    ("F&B Revenue",     "$0",           False, DGRAY),
    ("Grants",          "$0",           False, DGRAY),
    ("Total Revenue",   "$7,820",       True,  NAVY),
    ("",                "",             False, DGRAY),
    ("EXPENSE",         "",             True,  DGRAY),
    ("Perf. Fee",       "$2,100",       False, DGRAY),
    ("Tech Support",    "$240",         False, DGRAY),
    ("Total Expense",   "$2,340",       True,  NAVY),
    ("",                "",             False, DGRAY),
    ("Gross Margin",    "+$5,480",      True,  GREEN),
    ("Tickets Sold",    "148",          False, DGRAY),
]
for di, (k, v, bold, vc) in enumerate(detail_rows):
    dy = Inches(1.98 + di * 0.23)
    if k == "REVENUE" or k == "EXPENSE":
        txt(sl, k, Inches(0.55), dy, Inches(4.1), Inches(0.22),
            size=9, bold=True, color=DGRAY)
        continue
    if not k:
        continue
    txt(sl, k, Inches(0.55), dy, Inches(2.5), Inches(0.22),
        size=11, bold=bold, color=DGRAY if not bold else NAVY)
    txt(sl, v, Inches(3.0),  dy, Inches(1.7), Inches(0.22),
        size=11, bold=bold, color=vc, align=PP_ALIGN.RIGHT)

# edit panel
rect(sl, Inches(5.1), Inches(1.4), Inches(4.5), Inches(5.9), fill=LGRAY)
rect(sl, Inches(5.1), Inches(1.4), Inches(4.5), Inches(0.45), fill=GOLD)
txt(sl, "Edit Event",
    Inches(5.25), Inches(1.44), Inches(4.2), Inches(0.38),
    size=13, bold=True, color=NAVY)

edit_fields = [
    ("Date",             "2026-04-26"),
    ("Event Name",       "Kevin Hays Trio"),
    ("Series",           "Performance  ▾"),
    ("Program",          "Performance  ▾"),
    ("Artists",          "Kevin Hays  [+ Add]"),
    ("Ticket Sales",     "7820"),
    ("Perf. Fee",        "2100"),
    ("Tech Support",     "240"),
    ("Tickets Sold",     "148"),
    ("Notes",            ""),
]
for fi, (lbl, val) in enumerate(edit_fields):
    fy = Inches(2.0 + fi * 0.5)
    txt(sl, lbl, Inches(5.22), fy, Inches(1.6), Inches(0.26),
        size=10, color=DGRAY)
    rect(sl, Inches(6.85), fy - Inches(0.02), Inches(2.55), Inches(0.3),
         fill=WHITE)
    txt(sl, val, Inches(6.92), fy, Inches(2.4), Inches(0.28),
        size=11, color=NAVY)

rect(sl, Inches(7.3), Inches(7.08), Inches(1.8), Inches(0.35), fill=NAVY)
txt(sl, "Save Changes", Inches(7.3), Inches(7.1), Inches(1.8), Inches(0.3),
    size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# duplicate callout
rect(sl, Inches(9.9), Inches(1.4), Inches(3.05), Inches(5.9), fill=NAVY2)
txt(sl, "What you can do\nfrom this screen",
    Inches(10.05), Inches(1.55), Inches(2.75), Inches(0.75),
    size=15, bold=True, color=GOLD2)

actions = [
    ("View full P&L",       "Every revenue and expense\nline for the event."),
    ("Edit any field",      "Change dates, financials,\nartists, notes — anything."),
    ("Duplicate",           "Clone a recurring event\nas a starting point."),
    ("Delete",              "Remove with a single\nconfirm click."),
    ("Artist drill-down",   "Click an artist name to\nsee their full history."),
]
for ai, (head, body) in enumerate(actions):
    ay = Inches(2.4 + ai * 1.0)
    txt(sl, "▸ " + head, Inches(10.05), ay, Inches(2.75), Inches(0.3),
        size=13, bold=True, color=GOLD2)
    txt(sl, body, Inches(10.2), ay + Inches(0.3), Inches(2.6), Inches(0.55),
        size=12, color=RGBColor(0xC5, 0xD5, 0xEA), wrap=True)


# ═══════════════════════════════════════════════════════════════════════════════
# Slide 7 — Feature: Artist Directory
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
rect(sl, 0, 0, W, Inches(1.25), fill=NAVY)
txt(sl, "Feature 3 of 4  —  Artist Directory", Inches(0.5), Inches(0.12),
    Inches(10), Inches(0.5), size=14, color=GOLD)
txt(sl, "Know what every artist earns across every show.", Inches(0.5), Inches(0.52),
    Inches(12), Inches(0.65), size=32, bold=True, color=WHITE)
divider(sl, Inches(1.25))

# artist list
rect(sl, Inches(0.4), Inches(1.4), Inches(5.5), Inches(5.9), fill=LGRAY)
rect(sl, Inches(0.4), Inches(1.4), Inches(5.5), Inches(0.42), fill=NAVY)
for lbl, xp in [("Artist", 0.55), ("Shows", 3.7), ("Revenue", 4.6)]:
    txt(sl, lbl, Inches(xp), Inches(1.44), Inches(1.5), Inches(0.34),
        size=11, bold=True, color=WHITE)

artist_rows = [
    ("Kevin Hays",           "14", "$48,320"),
    ("Endea Owens",          "9",  "$38,740"),
    ("Moe Flavour",          "11", "$31,200"),
    ("Francine Reed",        "8",  "$27,880"),
    ("Jocelyn Gould Quartet","6",  "$22,140"),
    ("Sonrisa",              "7",  "$18,990"),
    ("Dave Potter",          "12", "$17,650"),
    ("Marcus Printup",       "5",  "$15,400"),
    ("Various (Education)",  "22", "$14,300"),
]
for ri, (name, shows, rev) in enumerate(artist_rows):
    ry = Inches(1.88 + ri * 0.52)
    rect(sl, Inches(0.4), ry, Inches(5.5), Inches(0.5),
         fill=WHITE if ri % 2 == 0 else LGRAY)
    txt(sl, name,  Inches(0.55), ry + Inches(0.1), Inches(3.0), Inches(0.3),
        size=12, color=NAVY)
    txt(sl, shows, Inches(3.7),  ry + Inches(0.1), Inches(0.8), Inches(0.3),
        size=12, color=DGRAY, align=PP_ALIGN.CENTER)
    txt(sl, rev,   Inches(4.55), ry + Inches(0.1), Inches(1.2), Inches(0.3),
        size=12, bold=True, color=GREEN, align=PP_ALIGN.RIGHT)

# artist detail
rect(sl, Inches(6.2), Inches(1.4), Inches(6.75), Inches(5.9), fill=LGRAY)
rect(sl, Inches(6.2), Inches(1.4), Inches(6.75), Inches(0.42), fill=GOLD)
txt(sl, "Kevin Hays  —  14 shows  ·  $48,320 revenue",
    Inches(6.35), Inches(1.44), Inches(6.4), Inches(0.36),
    size=12, bold=True, color=NAVY)

rect(sl, Inches(6.2), Inches(1.88), Inches(6.75), Inches(0.36), fill=NAVY)
for lbl, xp in [("Date", 6.35), ("Show", 7.5), ("Share", 9.5), ("Revenue", 10.5), ("Margin", 11.6)]:
    txt(sl, lbl, Inches(xp), Inches(1.91), Inches(1.1), Inches(0.28),
        size=10, bold=True, color=WHITE)

detail_rows2 = [
    ("Apr 26 '26", "Kevin Hays Trio",   "100%", "$7,820", "+$5,480"),
    ("Jan 18 '26", "KH Solo Evening",   "100%", "$5,200", "+$3,200"),
    ("Sep 12 '25", "Kevin Hays Trio",   "100%", "$6,640", "+$4,910"),
    ("Jun  7 '25", "KH Trio + Support", "70%",  "$4,830", "+$3,400"),
    ("Mar 15 '25", "Kevin Hays Quartet","100%", "$5,810", "+$3,980"),
    ("Nov  8 '24", "KH Solo",           "100%", "$4,200", "+$2,750"),
    ("Aug  3 '24", "KH Trio",           "100%", "$5,910", "+$4,100"),
]
for ri2, (date, show, share, rev, margin) in enumerate(detail_rows2):
    ry2 = Inches(2.3 + ri2 * 0.52)
    rect(sl, Inches(6.2), ry2, Inches(6.75), Inches(0.5),
         fill=WHITE if ri2 % 2 == 0 else LGRAY)
    for val, xp, wid, align, col in [
        (date,   6.35, 1.1, PP_ALIGN.LEFT,  NAVY),
        (show,   7.5,  1.9, PP_ALIGN.LEFT,  NAVY),
        (share,  9.5,  0.95,PP_ALIGN.CENTER,DGRAY),
        (rev,    10.5, 1.0, PP_ALIGN.RIGHT, DGRAY),
        (margin, 11.5, 1.3, PP_ALIGN.RIGHT, GREEN),
    ]:
        txt(sl, val, Inches(xp), ry2 + Inches(0.1),
            Inches(wid), Inches(0.3), size=11, color=col, align=align)

rect(sl, Inches(6.2), Inches(7.05), Inches(6.75), Inches(0.35), fill=NAVY2)
txt(sl, "14 shows  ·  $48,320 revenue  ·  $33,820 expense  ·  +$14,500 margin",
    Inches(6.35), Inches(7.08), Inches(6.4), Inches(0.28),
    size=11, bold=True, color=GOLD2)

txt(sl, "Share % lets you split revenue when an artist plays part of a multi-artist bill.",
    Inches(0.4), Inches(7.32), Inches(12.5), Inches(0.35),
    size=13, italic=True, color=DGRAY, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# Slide 8 — Feature: Export & Add New
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
rect(sl, 0, 0, W, Inches(1.25), fill=NAVY)
txt(sl, "Feature 4 of 4  —  Export & Add New Events", Inches(0.5), Inches(0.12),
    Inches(10), Inches(0.5), size=14, color=GOLD)
txt(sl, "Data in. Data out. Always in sync.", Inches(0.5), Inches(0.52),
    Inches(12), Inches(0.65), size=32, bold=True, color=WHITE)
divider(sl, Inches(1.25))

# export card
rect(sl, Inches(0.4), Inches(1.4), Inches(6.1), Inches(5.9), fill=LGRAY)
rect(sl, Inches(0.4), Inches(1.4), Inches(6.1), Inches(0.44), fill=GREEN)
txt(sl, "Export to Excel — anytime", Inches(0.55), Inches(1.44),
    Inches(5.8), Inches(0.36), size=15, bold=True, color=WHITE)

rect(sl, Inches(0.5), Inches(1.97), Inches(5.9), Inches(0.7),
     fill=RGBColor(0x1E, 0x1E, 0x1E))
txt(sl, "python export_to_excel.py",
    Inches(0.65), Inches(2.02), Inches(5.5), Inches(0.3),
    size=13, bold=True, color=RGBColor(0x9C, 0xD5, 0xFF))
txt(sl, "→  events_export.xlsx  (579 events)",
    Inches(0.65), Inches(2.33), Inches(5.5), Inches(0.28),
    size=12, color=RGBColor(0x6A, 0x99, 0x55))

export_points = [
    "·  One row per event — every financial field",
    "·  Total Revenue & Total Expense are live SUM()\n   formulas — edit a cell, totals update",
    "·  Gross Margin = Revenue − Expense (live formula)",
    "·  Revenue columns shaded green, expense red",
    "·  Artist names, ticket counts, and notes included",
    "·  Works as a backup — full data snapshot",
]
multiline(sl, export_points,
          Inches(0.55), Inches(2.82), Inches(5.7), Inches(4.0),
          size=13, color=NAVY, line_space=4)

# new event card
rect(sl, Inches(6.9), Inches(1.4), Inches(6.1), Inches(5.9), fill=LGRAY)
rect(sl, Inches(6.9), Inches(1.4), Inches(6.1), Inches(0.44), fill=GOLD)
txt(sl, "Add a new event in the browser", Inches(7.05), Inches(1.44),
    Inches(5.8), Inches(0.36), size=15, bold=True, color=NAVY)

# mini new-event form
new_fields = [
    ("Date",        "2026-05-03"),
    ("Event Name",  "Trio Night May"),
    ("Series",      "Performance  ▾"),
    ("Program",     "Performance  ▾"),
    ("Artist",      "Kevin Hays  [+ Add]"),
    ("Ticket Sales",""),
    ("Perf. Fee",   ""),
    ("Tickets Sold",""),
]
for fi, (lbl, val) in enumerate(new_fields):
    fy2 = Inches(1.98 + fi * 0.52)
    txt(sl, lbl, Inches(7.05), fy2, Inches(1.6), Inches(0.28),
        size=11, color=DGRAY)
    rect(sl, Inches(8.7), fy2 - Inches(0.02), Inches(3.9), Inches(0.32),
         fill=WHITE)
    if val:
        txt(sl, val, Inches(8.78), fy2, Inches(3.75), Inches(0.28),
            size=11, color=NAVY)

rect(sl, Inches(11.0), Inches(6.65), Inches(1.8), Inches(0.38), fill=NAVY)
txt(sl, "Save Event", Inches(11.0), Inches(6.67), Inches(1.8), Inches(0.32),
    size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

txt(sl, "Fill in what you know now — every field except Date and Name is optional.",
    Inches(0.4), Inches(7.32), Inches(12.5), Inches(0.35),
    size=13, italic=True, color=DGRAY, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# Slide 9 — Getting Started
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
rect(sl, 0, 0, W, Inches(1.25), fill=NAVY)
txt(sl, "Getting Started", Inches(0.5), Inches(0.18),
    Inches(10), Inches(0.65), size=34, bold=True, color=WHITE)
txt(sl, "From zero to running in three steps.",
    Inches(0.5), Inches(0.78), Inches(12.5), Inches(0.4),
    size=16, color=GOLD2)
divider(sl, Inches(1.25))

steps = [
    ("1", "Load the data",
     "Run the one-time import script against\nthe existing Excel files.\n\nAll historical events from FY2022\nthrough today are loaded automatically.\nNo manual entry required.",
     NAVY),
    ("2", "Start the app",
     "Launch Nash Analysis on a laptop\nor shared server.\n\nOpen a browser. Log in.\nYou're looking at all 579 events,\nsearchable and sortable.",
     GREEN),
    ("3", "Use it",
     "Add new events after each show.\nEdit financials as actuals come in.\nSearch for any show in seconds.\nExport to Excel whenever you need it.",
     GOLD),
]

for si, (num, title, body, color) in enumerate(steps):
    sx = Inches(0.5 + si * 4.25)
    rect(sl, sx, Inches(1.5), Inches(3.95), Inches(5.5), fill=LGRAY)
    rect(sl, sx, Inches(1.5), Inches(3.95), Inches(0.55), fill=color)
    txt(sl, num, sx + Inches(0.12), Inches(1.56),
        Inches(0.38), Inches(0.48), size=30, bold=True,
        color=WHITE if color != GOLD else NAVY)
    txt(sl, title, sx + Inches(0.72), Inches(1.58),
        Inches(3.1), Inches(0.44), size=18, bold=True, color=NAVY)
    txt(sl, body,  sx + Inches(0.72), Inches(2.15),
        Inches(3.1), Inches(4.6), size=14, color=DGRAY, wrap=True)

txt(sl, "No new software to license. No training on a foreign system. It works like the web — because it is.",
    Inches(0.4), Inches(7.1), Inches(12.5), Inches(0.35),
    size=14, italic=True, color=DGRAY, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# Slide 10 — Close / Call to Action
# ═══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
rect(sl, 0, 0, W, H, fill=NAVY)
rect(sl, 0, 0, Inches(0.55), H, fill=GOLD)
rect(sl, 0, Inches(4.1), W, Inches(0.07), fill=GOLD)

txt(sl, "Stop searching.\nStart knowing.",
    Inches(0.9), Inches(0.55), Inches(11.5), Inches(2.5),
    size=52, bold=True, color=WHITE)

txt(sl, "Nash Analysis puts every show's financials at your fingertips —\nno file-hunting, no version confusion, no manual aggregation.",
    Inches(0.9), Inches(3.1), Inches(11.5), Inches(0.9),
    size=18, color=GOLD2)

ctas = [
    ("Live demo",        "See the app running against real data."),
    ("Pilot it",         "Run alongside existing Excel files for one month."),
    ("Full rollout",     "Load all historical data and go live for the team."),
]
for ci, (action, desc) in enumerate(ctas):
    cx = Inches(0.9 + ci * 4.1)
    rect(sl, cx, Inches(4.4), Inches(3.8), Inches(1.55), fill=NAVY2)
    rect(sl, cx, Inches(4.4), Inches(0.12), Inches(1.55), fill=GOLD)
    txt(sl, action, cx + Inches(0.25), Inches(4.52),
        Inches(3.45), Inches(0.42), size=16, bold=True, color=WHITE)
    txt(sl, desc,   cx + Inches(0.25), Inches(4.95),
        Inches(3.45), Inches(0.85), size=13,
        color=RGBColor(0xC5, 0xD5, 0xEA), wrap=True)

txt(sl, "What's the right next step?",
    Inches(0.9), Inches(6.1), Inches(11.5), Inches(0.45),
    size=20, bold=True, color=GOLD2)

txt(sl, "Nash Jazz Club  ·  Finance & Operations",
    Inches(0.9), Inches(7.05), Inches(11.5), Inches(0.38),
    size=12, italic=True, color=DGRAY)


prs.save(out_path)
print(f"Saved → {out_path}  ({len(prs.slides)} slides)")
